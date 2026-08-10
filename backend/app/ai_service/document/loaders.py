"""
文档格式加载器：md / docx / pdf 文字与内嵌资源提取。
"""

from __future__ import annotations

import logging
import re
import tempfile
from pathlib import Path

from .office_convert import convert_with_libreoffice, libreoffice_unavailable_message
from .schemas import DocumentBlock, DocumentBlockType

log = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".doc",
    ".docx",
    ".pdf",
    ".pptx",
    ".xlsx",
}


def _read_markdown(path: Path) -> tuple[str, list[Path]]:
    """读取 Markdown/纯文本，并收集相对路径图片引用。"""
    text = path.read_text(encoding="utf-8", errors="replace")
    image_paths: list[Path] = []
    base_dir = path.parent
    for match in re.finditer(r"!\[[^\]]*\]\(([^)]+)\)", text):
        ref = match.group(1).strip().split(" ")[0]
        if ref.startswith(("http://", "https://", "data:")):
            continue
        candidate = (base_dir / ref).resolve()
        if candidate.is_file():
            image_paths.append(candidate)
    return text, image_paths


def _read_docx(path: Path) -> tuple[str, list[Path]]:
    """读取 docx 段落文字并导出内嵌图片到临时目录。"""
    from docx import Document  # python-docx

    from .docx_compat import ensure_docx_broken_rels_patch

    ensure_docx_broken_rels_patch()

    try:
        doc = Document(str(path))
    except KeyError as exc:
        # 仍失败时尝试 LibreOffice 另存清洗后再读
        log.warning("docx open failed (%s), try LibreOffice repair: %s", path.name, exc)
        tmp_dir = Path(tempfile.mkdtemp(prefix="kb_docx_repair_"))
        repaired = convert_with_libreoffice(path, tmp_dir, "docx")
        if repaired is None:
            raise ValueError(
                f"无法打开 Word 文档（损坏的内部链接/图片引用: {exc}）。"
                "请用 Word 打开后另存为新的 .docx 再上传。"
            ) from exc
        doc = Document(str(repaired))

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # 表格文字一并纳入，避免「正文在表里」时切分结果为空
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    text = "\n\n".join(paragraphs)

    image_paths: list[Path] = []
    tmp_dir = Path(tempfile.mkdtemp(prefix="kb_docx_img_"))
    img_idx = 0
    for rel in doc.part.rels.values():
        if "image" not in getattr(rel, "reltype", ""):
            continue
        try:
            if getattr(rel, "is_external", False):
                continue
            target_ref = (getattr(rel, "target_ref", None) or "").strip()
            if not target_ref or target_ref.upper() == "NULL" or target_ref.startswith("#"):
                continue
            part = rel.target_part
            blob = part.blob
            ext = part.content_type.split("/")[-1].replace("jpeg", "jpg")
            out = tmp_dir / f"img_{img_idx}.{ext}"
            out.write_bytes(blob)
            image_paths.append(out)
            img_idx += 1
        except Exception as exc:  # noqa: BLE001
            log.warning("docx 内嵌图跳过 (%s): %s", path.name, exc)
    return text, image_paths


def _read_pdf(path: Path) -> tuple[str, list[Path]]:
    """读取 PDF 文字层；对无文字页导出整页图片供 VL 处理。"""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    text_parts: list[str] = []
    image_paths: list[Path] = []
    tmp_dir = Path(tempfile.mkdtemp(prefix="kb_pdf_img_"))

    for page_index in range(len(doc)):
        page = doc[page_index]
        page_text = page.get_text("text").strip()
        if page_text:
            text_parts.append(f"[第{page_index + 1}页]\n{page_text}")
        else:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            out = tmp_dir / f"page_{page_index + 1}.png"
            pix.save(str(out))
            image_paths.append(out)

        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base = doc.extract_image(xref)
                ext = base.get("ext", "png")
                img_bytes = base.get("image")
                if not img_bytes:
                    continue
                out = tmp_dir / f"page{page_index + 1}_img_{xref}.{ext}"
                out.write_bytes(img_bytes)
                image_paths.append(out)
            except Exception as exc:
                log.warning("PDF 内嵌图提取失败 xref=%s: %s", xref, exc)

    doc.close()
    return "\n\n".join(text_parts), image_paths


def _read_doc(path: Path) -> tuple[str, list[Path]]:
    """读取旧版 .doc：优先 LibreOffice 转 docx，否则提示安装。"""
    tmp_dir = Path(tempfile.mkdtemp(prefix="kb_doc_convert_"))
    converted = convert_with_libreoffice(path, tmp_dir, "docx")
    if converted is None:
        raise ValueError(libreoffice_unavailable_message())
    return _read_docx(converted)


def _read_pptx(path: Path) -> tuple[str, list[Path]]:
    """读取 pptx 幻灯片文字并导出内嵌图片。"""
    from pptx import Presentation

    prs = Presentation(str(path))
    text_parts: list[str] = []
    image_paths: list[Path] = []
    tmp_dir = Path(tempfile.mkdtemp(prefix="kb_pptx_img_"))
    img_idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_lines.append(shape.text.strip())
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    out = tmp_dir / f"slide{slide_num}_img_{img_idx}.{ext}"
                    out.write_bytes(image.blob)
                    image_paths.append(out)
                    img_idx += 1
                except Exception as exc:
                    log.warning("PPTX 图片提取失败 slide=%s: %s", slide_num, exc)
        if slide_lines:
            text_parts.append(f"[第{slide_num}页]\n" + "\n".join(slide_lines))

    return "\n\n".join(text_parts), image_paths


def _read_xlsx(path: Path) -> tuple[str, list[Path]]:
    """读取 xlsx 各 Sheet 单元格文字（只读模式）。"""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
    finally:
        wb.close()
    return "\n\n".join(parts), []


def load_document_text_and_images(path: Path) -> tuple[str, list[Path], list[str]]:
    """
    按后缀加载文档，返回 (纯文本, 图片路径列表, 警告信息)。

    :raises ValueError: 不支持的格式
    """
    suffix = path.suffix.lower()
    warnings: list[str] = []
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文档格式: {suffix}")

    if suffix in {".md", ".markdown", ".txt"}:
        text, images = _read_markdown(path)
    elif suffix == ".doc":
        text, images = _read_doc(path)
    elif suffix == ".docx":
        text, images = _read_docx(path)
    elif suffix == ".pdf":
        text, images = _read_pdf(path)
    elif suffix == ".pptx":
        text, images = _read_pptx(path)
    elif suffix == ".xlsx":
        text, images = _read_xlsx(path)
    else:
        raise ValueError(f"不支持的文档格式: {suffix}")

    if not text and not images:
        warnings.append("未提取到文字或图片内容")
    return text, images, warnings


def text_to_blocks(text: str, *, source: str) -> list[DocumentBlock]:
    """将整段文字转为 TEXT 类型 DocumentBlock。"""
    if not text.strip():
        return []
    return [
        DocumentBlock(
            block_type=DocumentBlockType.TEXT,
            text=text.strip(),
            source=source,
        )
    ]
